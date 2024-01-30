from pptx import Presentation
from io import BytesIO


def generate_presentation(text):

    # Splitting the summary text into topics
    topics = text.split('\n\n')

    # Creating presentation object
    presentation = Presentation()

    # Creating slide layout
    first_slide_layout = presentation.slide_layouts[0]

    # Creating slide object to add in ppt
    slide = presentation.slides.add_slide(first_slide_layout)
    slide.shapes.title.text = "pdf to ppt project"
    slide.placeholders[1].text = "By Tanvi, Toshal and Vaish"

    # Adding slides dynamically for each topic with bullet points
    for topic in topics:
        slide_layout = presentation.slide_layouts[1]  # Assuming you want title and content layout
        slide = presentation.slides.add_slide(slide_layout)
        title, points = topic.split('Points:')
        slide.shapes.title.text = title.strip()
        slide.placeholders[1].text = points.strip()

    # Saving the presentation to a BytesIO buffer
    ppt_buffer = BytesIO()
    presentation.save(ppt_buffer)
    ppt_buffer.seek(0)

    return ppt_buffer

